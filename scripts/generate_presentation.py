from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "AI_Code_Review_Assistant_Presentation.pptx"

BG = RGBColor(15, 23, 42)
ACCENT = RGBColor(56, 189, 248)
TEXT = RGBColor(241, 245, 249)
MUTED = RGBColor(148, 163, 184)
GOOD = RGBColor(74, 222, 128)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.2), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TEXT

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.25), Inches(1.2), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.color.rgb = ACCENT

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.0), Inches(0.5))
        p = subtitle_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(14)
        r.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], *, left: float = 0.9, top: float = 1.9, width: float = 11.0, height: float = 4.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT
        p.space_after = Pt(12)


def add_two_column(slide, left_title: str, left_items: list[str], right_title: str, right_items: list[str]) -> None:
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.3), Inches(4.8))
    p = left_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = left_title
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    for item in left_items:
        p = left_box.text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)

    right_box = slide.shapes.add_textbox(Inches(6.4), Inches(1.9), Inches(5.3), Inches(4.8))
    p = right_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = right_title
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    for item in right_items:
        p = right_box.text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(11.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "AI-Powered Code Review Assistant", "Manual review + GitHub PR review in one application")
    add_bullets(
        slide,
        [
            "Built as a learning project to review pasted code, diffs, and GitHub pull requests.",
            "Combines Streamlit UI, FastAPI webhook handling, LangChain review orchestration, and GitHub integration.",
            "Includes a local fallback reviewer so the product still works when external LLM access is unavailable.",
        ],
        top=2.0,
        height=3.0,
    )
    add_footer(slide, "Project presentation")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Problem Statement", "Why this project was built")
    add_bullets(
        slide,
        [
            "Developers need faster feedback on correctness, security, testing gaps, and maintainability risks.",
            "Manual review is time-consuming and inconsistent across diffs, code snippets, and pull requests.",
            "The goal was to create one app that supports both manual review and GitHub-based review workflows.",
        ],
    )
    add_footer(slide, "Problem and objective")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Solution Overview")
    add_two_column(
        slide,
        "Core Features",
        [
            "Manual review tab for pasted diffs, patches, or code",
            "GitHub PR review tab inside the same Streamlit app",
            "FastAPI webhook endpoint for PR-triggered review",
            "Review history stored in SQLite",
        ],
        "Supporting Features",
        [
            "Repository context retrieval for nearby code awareness",
            "Structured findings, risk, and missing tests output",
            "GitHub comment and status publication when permitted",
            "Fallback local reviewer when model calls fail",
        ],
    )
    add_footer(slide, "What the application delivers")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Architecture")
    add_bullets(
        slide,
        [
            "Streamlit app (`app.py`) handles manual reviews and GitHub PR review preview/publish flows.",
            "FastAPI service (`webhook_app.py`) receives GitHub webhook events and queues PR review processing.",
            "Review engine (`reviewer.py`) builds the LangChain prompt pipeline and structured output handling.",
            "GitHub service layer (`github_review.py`, `github_client.py`) fetches PR files and posts results.",
            "SQLite history store keeps recent review results for both manual and GitHub sources.",
        ],
        height=4.9,
    )
    add_footer(slide, "Application structure")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "How LangChain Was Used")
    add_bullets(
        slide,
        [
            "A system prompt defines the review behavior: focus on correctness, regressions, security, performance, and testing gaps.",
            "A user prompt injects repository context, retrieved files, changed file summary, and submitted code/diff.",
            "LangChain structured output maps the model response into a typed review schema.",
            "Primary files: `src/code_review_assistant/reviewer.py` and `src/code_review_assistant/models.py`.",
        ],
    )
    add_footer(slide, "LangChain workflow")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Key Files")
    add_two_column(
        slide,
        "Frontend / API",
        [
            "`app.py` - Streamlit UI",
            "`webhook_app.py` - FastAPI webhook service",
            "`README.md` - setup and demo guide",
        ],
        "Backend Logic",
        [
            "`reviewer.py` - LangChain review orchestration",
            "`github_review.py` - PR review workflow",
            "`repository_context.py` - repository retrieval",
            "`fallback_reviewer.py` - local offline review path",
        ],
    )
    add_footer(slide, "Important implementation files")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Fallback Review Mode")
    add_bullets(
        slide,
        [
            "Added because external Gemini model access was unavailable due to quota / billing constraints.",
            "The fallback reviewer handles obvious hardcoded secrets, debug statements, broad exceptions, likely missing semicolons, and missing test reminders.",
            "This keeps the app usable for demo and learning even without full LLM availability.",
            "It is narrower than a full model-based review and should be presented as a resilience feature, not a replacement for production-grade LLM review.",
        ],
        height=5.0,
    )
    add_footer(slide, "Known limitation handled pragmatically")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "How To Run")
    add_bullets(
        slide,
        [
            "1. `pip install -r requirements.txt`",
            "2. `copy .env.example .env`",
            "3. Update `.env` with provider and GitHub credentials",
            "4. `streamlit run app.py`",
            "5. Optional webhook mode: `uvicorn webhook_app:app --reload --port 8000`",
        ],
    )
    add_footer(slide, "Setup commands")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Demo Flow")
    add_bullets(
        slide,
        [
            "Manual Review: paste a code snippet or diff, set the primary language, and run review.",
            "GitHub PR Review: enter owner, repository, and PR number in the GitHub tab.",
            "Preview Mode: keep publish disabled to review results inside the app without posting to GitHub.",
            "Show review history at the bottom of the app to demonstrate stored results.",
        ],
    )
    add_footer(slide, "Suggested presentation demo")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Challenges And Limitations")
    add_two_column(
        slide,
        "Challenges",
        [
            "Gemini quota / billing was unavailable",
            "GitHub posting depends on token or app permissions",
            "Raw pasted code needs different handling than unified diffs",
        ],
        "Current Limitations",
        [
            "Fallback review is narrower than LLM review",
            "GitHub publish mode can fail with insufficient token scopes",
            "Full LLM quality requires enabled provider access",
        ],
    )
    add_footer(slide, "Transparent project status")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Testing And Validation")
    add_bullets(
        slide,
        [
            "Test coverage includes parser logic, GitHub review helpers, webhook verification, deduplication, repository context, and fallback review.",
            "Final validation status: all automated tests pass.",
            "The app supports manual review, GitHub PR preview review, and webhook-based review handling.",
        ],
    )
    add_footer(slide, "Quality checks")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "The project successfully delivers a unified code review assistant for both manual and GitHub PR workflows.",
            "LangChain was used to structure the review pipeline and typed output flow.",
            "A fallback review engine was added to keep the application usable under provider limitations.",
            "The project is runnable, tested, and ready for demonstration.",
        ],
    )
    add_footer(slide, "Final outcome")

    return prs


def main() -> int:
    presentation = build_presentation()
    presentation.save(OUTPUT)
    print(f"Created {OUTPUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
